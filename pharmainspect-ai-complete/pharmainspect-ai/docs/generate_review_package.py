from __future__ import annotations

import re
from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocInches
from docx.shared import Pt, RGBColor as DocRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parent
SOURCE = ROOT / "PharmaInspect_AI_Architecture_and_Business_Review.md"
DOCX_OUTPUT = ROOT / "PharmaInspect_AI_Architecture_and_Business_Review.docx"
PPTX_OUTPUT = ROOT / "PharmaInspect_AI_Executive_Review.pptx"

NAVY = RGBColor(7, 57, 87)
DEEP_NAVY = RGBColor(5, 38, 61)
BLUE = RGBColor(14, 165, 233)
GREEN = RGBColor(16, 185, 129)
AMBER = RGBColor(245, 158, 11)
RED = RGBColor(239, 68, 68)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
PALE_BLUE = RGBColor(240, 249, 255)
PALE_GREEN = RGBColor(236, 253, 245)
PALE_AMBER = RGBColor(255, 251, 235)
PALE_RED = RGBColor(254, 242, 242)
WHITE = RGBColor(255, 255, 255)
LIGHT_BORDER = RGBColor(224, 242, 254)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill)
    tc_pr.append(shading)


def build_docx() -> None:
    document = Document()
    section = document.sections[0]
    section.top_margin = DocInches(0.65)
    section.bottom_margin = DocInches(0.65)
    section.left_margin = DocInches(0.75)
    section.right_margin = DocInches(0.75)

    styles = document.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(9.5)
    for level in range(1, 4):
        style = styles[f"Heading {level}"]
        style.font.name = "Aptos Display"
        style.font.color.rgb = DocRGBColor(7, 89, 133)

    lines = SOURCE.read_text(encoding="utf-8").splitlines()
    index = 0
    in_code = False
    code_lines: list[str] = []

    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                paragraph = document.add_paragraph()
                run = paragraph.add_run("\n".join(code_lines))
                run.font.name = "Consolas"
                run.font.size = Pt(8)
                paragraph.paragraph_format.space_after = Pt(8)
                code_lines = []
                in_code = False
            else:
                in_code = True
            index += 1
            continue

        if in_code:
            code_lines.append(line)
            index += 1
            continue

        if not stripped or stripped == "---":
            index += 1
            continue

        if stripped.startswith("|") and index + 1 < len(lines) and re.match(
            r"^\s*\|(?:\s*:?-+:?\s*\|)+\s*$", lines[index + 1]
        ):
            table_lines = [stripped]
            index += 2
            while index < len(lines) and lines[index].strip().startswith("|"):
                table_lines.append(lines[index].strip())
                index += 1
            rows = [
                [cell.strip() for cell in row.strip("|").split("|")]
                for row in table_lines
            ]
            table = document.add_table(rows=1, cols=len(rows[0]))
            table.style = "Table Grid"
            for column, value in enumerate(rows[0]):
                table.rows[0].cells[column].text = value
                set_cell_shading(table.rows[0].cells[column], "E0F2FE")
                for run in table.rows[0].cells[column].paragraphs[0].runs:
                    run.bold = True
                    run.font.color.rgb = DocRGBColor(7, 89, 133)
            for row_values in rows[1:]:
                cells = table.add_row().cells
                for column, value in enumerate(row_values):
                    cells[column].text = value
            document.add_paragraph()
            continue

        heading = re.match(r"^(#{1,4})\s+(.+)$", stripped)
        if heading:
            level = min(len(heading.group(1)), 3)
            paragraph = document.add_heading(heading.group(2), level=level)
            if level == 1:
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            index += 1
            continue

        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:], style="List Bullet")
        elif re.match(r"^\d+\.\s+", stripped):
            document.add_paragraph(re.sub(r"^\d+\.\s+", "", stripped), style="List Number")
        elif stripped.startswith("**") and stripped.endswith("**"):
            paragraph = document.add_paragraph()
            paragraph.add_run(stripped.strip("*")).bold = True
        else:
            paragraph = document.add_paragraph(stripped)
            paragraph.paragraph_format.space_after = Pt(5)
        index += 1

    document.add_section(WD_SECTION.NEW_PAGE)
    paragraph = document.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run("End of architecture and business review")
    run.italic = True
    run.font.color.rgb = DocRGBColor(100, 116, 139)
    document.save(DOCX_OUTPUT)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 18,
    color: RGBColor = NAVY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, subtitle: str | None = None, section: str | None = None) -> None:
    if section:
        add_text(slide, section.upper(), 0.65, 0.35, 4.5, 0.3, size=9, color=BLUE, bold=True)
    add_text(slide, title, 0.65, 0.72, 12.0, 0.58, size=27, color=NAVY, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.67, 1.34, 11.8, 0.45, size=12, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.67), Inches(1.83), Inches(12), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BORDER
    line.line.fill.background()


def add_footer(slide, number: int) -> None:
    add_text(slide, "PharmaInspect AI | Architecture and Business Review", 0.65, 7.14, 6.6, 0.2, size=7, color=MUTED)
    add_text(slide, str(number), 12.25, 7.12, 0.4, 0.22, size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(
    slide,
    x: float,
    y: float,
    width: float,
    height: float,
    title: str,
    body: str,
    *,
    fill: RGBColor = WHITE,
    accent: RGBColor = BLUE,
    title_size: float = 14,
    body_size: float = 10.5,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LIGHT_BORDER
    accent_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(0.07),
        Inches(height),
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = accent
    accent_shape.line.fill.background()
    if height < 0.85:
        title_width = min(1.7, width * 0.34)
        add_text(
            slide,
            title,
            x + 0.2,
            y + 0.16,
            title_width,
            height - 0.22,
            size=title_size,
            color=NAVY,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            body,
            x + 0.28 + title_width,
            y + 0.16,
            width - title_width - 0.48,
            height - 0.22,
            size=body_size,
            color=SLATE,
            valign=MSO_ANCHOR.MIDDLE,
        )
        return shape
    add_text(slide, title, x + 0.22, y + 0.18, width - 0.4, 0.32, size=title_size, color=NAVY, bold=True)
    add_text(slide, body, x + 0.22, y + 0.62, width - 0.4, height - 0.75, size=body_size, color=SLATE)
    return shape


def add_bullets(
    slide,
    items: list[str],
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float = 14,
    color: RGBColor = SLATE,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.05)
    for item_index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = PptPt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = PptPt(10)
        paragraph.text = f"•  {item}"
    return box


def add_metric(slide, x: float, y: float, value: str, label: str, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.35), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LIGHT_BORDER
    add_text(slide, value, x + 0.15, y + 0.15, 2.05, 0.45, size=24, color=color, bold=True)
    add_text(slide, label, x + 0.15, y + 0.68, 2.05, 0.25, size=9, color=MUTED, bold=True)


def new_slide(presentation: Presentation, background: RGBColor = PALE_BLUE):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background_fill = slide.background.fill
    background_fill.solid()
    background_fill.fore_color.rgb = background
    return slide


def build_pptx() -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # 1. Title
    slide = new_slide(presentation, DEEP_NAVY)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    slide.shapes[-1].line.fill.background()
    add_text(slide, "PHARMAINSPECT AI", 0.85, 0.72, 4.0, 0.35, size=12, color=GREEN, bold=True)
    add_text(slide, "Architecture and\nBusiness Review", 0.82, 1.45, 8.2, 1.6, size=38, color=WHITE, bold=True)
    add_text(
        slide,
        "Human-in-the-loop pharmaceutical packaging inspection\nwith retrieval-grounded visual evidence",
        0.86,
        3.35,
        7.8,
        0.9,
        size=17,
        color=RGBColor(186, 230, 253),
    )
    add_metric(slide, 0.86, 5.2, "50%", "RELEVANCE GATE", BLUE)
    add_metric(slide, 3.43, 5.2, "90%", "ACCEPTANCE GATE", GREEN)
    add_metric(slide, 6.0, 5.2, "Human", "FINAL AUTHORITY", AMBER)
    add_text(slide, "Current-state review and production roadmap | July 2026", 0.88, 6.83, 7.0, 0.25, size=9, color=MUTED)

    # 2. Executive summary
    slide = new_slide(presentation)
    add_title(slide, "Executive summary", "A complete prototype workflow with a clear production-readiness path", "Business review")
    add_card(slide, 0.7, 2.15, 3.85, 3.7, "What exists today", "Batch creation and image upload\n\nGood-image vector retrieval\n\nInvalid / fail / pass separation\n\nExplainable image evidence\n\nHuman review and PDF reporting", fill=WHITE, accent=BLUE)
    add_card(slide, 4.75, 2.15, 3.85, 3.7, "Business value", "Focus inspectors on exceptions\n\nStandardize first-pass screening\n\nPreserve image-to-disposition traceability\n\nReduce generic AI explanations\n\nCreate reusable structured evidence", fill=WHITE, accent=GREEN)
    add_card(slide, 8.8, 2.15, 3.85, 3.7, "Recommendation", "Proceed with a controlled pilot\n\nMeasure false accepts and rejects\n\nAdd persistent governed data\n\nValidate model and references\n\nComplete security and GxP controls", fill=WHITE, accent=AMBER)
    add_footer(slide, 2)

    # 3. Business problem
    slide = new_slide(presentation)
    add_title(slide, "Business problem and opportunity", "Inspection support must improve consistency without transferring release authority to AI", "Business review")
    add_card(slide, 0.7, 2.2, 5.8, 3.9, "Current challenges", "• Repetitive visual screening and reviewer fatigue\n• Evidence scattered across images, notes, and reports\n• Invalid uploads can distort quality metrics\n• Generic failure wording reduces trust\n• Batch decisions require attributable human review", fill=PALE_RED, accent=RED)
    add_card(slide, 6.8, 2.2, 5.8, 3.9, "Opportunity", "• Consistent first-pass screening\n• Retrieval against approved good examples\n• Explainable image-level evidence\n• Prioritized human review queue\n• Faster report preparation\n• Structured data for trends and CAPA workflows", fill=PALE_GREEN, accent=GREEN)
    add_footer(slide, 3)

    # 4. Scope
    slide = new_slide(presentation)
    add_title(slide, "Current solution scope", "Implemented workflow, explicit controls, and deliberate boundaries", "Product")
    add_card(slide, 0.7, 2.1, 3.8, 4.25, "Implemented", "Batch and image workflow\n\nReference-vector retrieval\n\nFeature-level explanations\n\nScoped quality assistant\n\nHuman review and disposition\n\nDashboard and PDF report", accent=GREEN)
    add_card(slide, 4.75, 2.1, 3.8, 4.25, "Safety controls", "50% relevance gate\n\n90% acceptance gate\n\nInvalid images excluded from KPIs\n\nUnknown entities are not guessed\n\nAssistant scope refusal\n\nHuman decision remains authoritative", accent=BLUE)
    add_card(slide, 8.8, 2.1, 3.8, 4.25, "Not yet production-ready", "No persistent database\n\nNo enterprise identity or RBAC\n\nNo immutable audit trail\n\nNo e-signature\n\nNo validated defect model\n\nNo HA / DR or formal CSV package", accent=AMBER)
    add_footer(slide, 4)

    # 5. Current architecture
    slide = new_slide(presentation)
    add_title(slide, "Current architecture", "Ports-and-adapters structure keeps vision, storage, reporting, and assistant components replaceable", "Architecture")
    layers = [
        ("React Quality Workspace", "Dashboard • Batch creation • Evidence • Review • Assistant", BLUE),
        ("FastAPI Application", "REST routes • Batch service • Metrics • Scope guard • Report service", NAVY),
        ("Infrastructure Adapters", "Image RAG • Vector index • Repository • Assistant • PDF", GREEN),
        ("Local Resources", "Approved image directory • In-memory batch state", AMBER),
    ]
    y = 2.12
    for title, body, color in layers:
        add_card(slide, 1.25, y, 10.85, 0.95, title, body, fill=WHITE, accent=color, title_size=13, body_size=9.5)
        y += 1.13
    add_text(slide, "Replaceable boundaries", 9.45, 6.72, 2.45, 0.25, size=9, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 5)

    # 6. Workflow
    slide = new_slide(presentation)
    add_title(slide, "End-to-end inspection workflow", "One evidence chain from batch context to human disposition", "Architecture")
    steps = [
        ("1", "Create batch", BLUE),
        ("2", "Upload images", BLUE),
        ("3", "Vectorize query", NAVY),
        ("4", "Retrieve good refs", GREEN),
        ("5", "Classify outcome", AMBER),
        ("6", "Human review", RED),
    ]
    x = 0.55
    for index, (number, label, color) in enumerate(steps):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.55), Inches(0.68), Inches(0.68))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, number, x, 2.67, 0.68, 0.25, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x - 0.28, 3.42, 1.25, 0.6, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        if index < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 0.86), Inches(2.67), Inches(0.88), Inches(0.42))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_BORDER
            arrow.line.fill.background()
        x += 2.08
    add_card(slide, 1.05, 4.75, 11.2, 1.15, "Authoritative record", "Structured result → batch metrics → inspector notes → root cause → corrective action → approve / hold / reject", fill=WHITE, accent=GREEN)
    add_footer(slide, 6)

    # 7. RAG
    slide = new_slide(presentation)
    add_title(slide, "Image RAG and decision thresholds", "Retrieval-grounded one-class anomaly detection using approved good-image evidence", "AI architecture")
    add_card(slide, 0.7, 2.08, 3.0, 2.2, "Vector composition", "35% color histogram\n40% luminance structure\n25% edge energy\n\nNormalized cosine similarity", accent=BLUE)
    add_card(slide, 3.95, 2.08, 2.65, 2.2, "Invalid", "< 50% relevance\n\nRejected before defect assessment\nExcluded from quality KPIs", fill=PALE_AMBER, accent=AMBER)
    add_card(slide, 6.85, 2.08, 2.65, 2.2, "Failed", "50% to < 90%\n\nValid packaging image\nHuman review required", fill=PALE_RED, accent=RED)
    add_card(slide, 9.75, 2.08, 2.65, 2.2, "Passed", "≥ 90% similarity\n\nConsistent with approved corpus\nHuman authority retained", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 0.7, 4.65, 11.7, 1.25, "Production interpretation", "Current retrieval is explainable and local. Production requires SKU-specific calibration, labeled defect data, governed reference versions, and a validated vision model.", fill=WHITE, accent=NAVY)
    add_footer(slide, 7)

    # 8. Explainability
    slide = new_slide(presentation)
    add_title(slide, "Explainability and invalid-image controls", "The system reports measurable reasons and protects metric integrity", "Controls")
    add_metric(slide, 0.75, 2.15, "25%", "EXAMPLE RANDOM IMAGE", RED)
    add_metric(slide, 3.35, 2.15, "50%", "MINIMUM RELEVANCE", AMBER)
    add_metric(slide, 5.95, 2.15, "90%", "PASS THRESHOLD", GREEN)
    add_metric(slide, 8.55, 2.15, "0", "KPI IMPACT IF INVALID", BLUE)
    add_card(slide, 0.75, 3.75, 5.7, 2.1, "Failure explanation", "Overall threshold shortfall\nPackage or label structural mismatch\nPackage or seal edge mismatch\nPackaging or label color mismatch", accent=RED)
    add_card(slide, 6.75, 3.75, 5.7, 2.1, "Invalid-image behavior", "Distinct “Invalid image” status\nNo defect assessment\nExcluded from pass/fail and defect counts\nAssistant explains the exclusion", accent=AMBER)
    add_footer(slide, 8)

    # 9. Assistant
    slide = new_slide(presentation)
    add_title(slide, "Grounded AI Quality Assistant", "Structured retrieval and deterministic guardrails prevent generic or unrelated answers", "AI architecture")
    add_card(slide, 0.75, 2.1, 3.7, 3.9, "Grounding", "Stored batch records\nExact image inspection results\nQuality metrics\nHuman review fields\nNamed batch and filename resolution", accent=BLUE)
    add_card(slide, 4.8, 2.1, 3.7, 3.9, "Guardrails", "Pharma-inspection scope check\nOut-of-scope refusal\nUnknown entity response\nRequested-image-only answer\nNo hidden reference filenames in response", accent=GREEN)
    add_card(slide, 8.85, 2.1, 3.7, 3.9, "TCS model integration", "TCS GenAI Lab endpoint\nDeepSeek V3 chat model\nStructured JSON grounding\nEnvironment-only API key\nTransparent local fallback", accent=AMBER)
    add_footer(slide, 9)

    # 10. Human and compliance
    slide = new_slide(presentation)
    add_title(slide, "Human review and compliance posture", "The workflow supports quality decisions; it is not yet a validated release system", "Governance")
    add_card(slide, 0.7, 2.1, 5.75, 3.9, "Controls already present", "• Human disposition remains authoritative\n• Structured notes, cause, and corrective action\n• Explicit invalid-image handling\n• Explainable thresholds and scores\n• Input MIME and size controls\n• Automated regression tests", fill=PALE_GREEN, accent=GREEN)
    add_card(slide, 6.8, 2.1, 5.75, 3.9, "Required before production", "• SSO, RBAC, and attributable identity\n• Immutable audit trail and e-signature\n• Persistent data and retention controls\n• Model/reference lifecycle governance\n• Security, performance, and recovery testing\n• Formal CSV and traceability package", fill=PALE_AMBER, accent=AMBER)
    add_footer(slide, 10)

    # 11. KPIs
    slide = new_slide(presentation)
    add_title(slide, "KPI and value framework", "Pilot value must be measured against an agreed operational baseline", "Business case")
    add_card(slide, 0.7, 2.1, 3.8, 3.75, "Quality KPIs", "False accept rate\nFalse reject rate\nInvalid-image precision\nReviewer override rate\nConfidence calibration\nDefect precision / recall", accent=RED)
    add_card(slide, 4.75, 2.1, 3.8, 3.75, "Operational KPIs", "Seconds saved per image\nReview time per batch\nEscalation rate\nBatches per reviewer\nUpload-to-disposition time\nReport preparation time", accent=BLUE)
    add_card(slide, 8.8, 2.1, 3.8, 3.75, "Value model", "Labor capacity benefit\n+\nAvoided quality-event value\n−\nPlatform and operating cost\n=\nNet annual value", accent=GREEN)
    add_text(slide, "No ROI claim should be approved until baseline and pilot measurements are available.", 1.3, 6.32, 10.7, 0.35, size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12. Risks
    slide = new_slide(presentation)
    add_title(slide, "Key risks and production gaps", "The largest risks are data quality, governance, and operational controls—not UI completion", "Risk review")
    risks = [
        ("Good-only corpus", "May miss defect diversity", "Add labeled bad examples", RED),
        ("Imaging variation", "Similarity drift", "Control camera and lighting", AMBER),
        ("Global thresholds", "Uneven SKU accuracy", "Calibrate per product/line", AMBER),
        ("In-memory state", "Evidence lost on restart", "Persistent database/storage", RED),
        ("Uncontrolled refs", "Untraceable decisions", "Versioned approval workflow", RED),
        ("No audit/RBAC", "Compliance exposure", "Identity, roles, immutable audit", RED),
    ]
    y = 2.05
    for title, impact, mitigation, color in risks:
        add_card(slide, 0.75, y, 3.3, 0.68, title, impact, accent=color, title_size=11, body_size=8.5)
        add_card(slide, 4.25, y, 8.25, 0.68, "Mitigation", mitigation, accent=GREEN, title_size=10, body_size=8.5)
        y += 0.76
    add_footer(slide, 12)

    # 13. Target architecture
    slide = new_slide(presentation)
    add_title(slide, "Production target architecture", "Durable, governed, observable, and horizontally scalable", "Architecture")
    rows = [
        [("Enterprise identity", BLUE), ("Gateway / WAF", BLUE), ("React web app", BLUE)],
        [("Inspection API", NAVY), ("Analysis queue", NAVY), ("Vision workers", NAVY)],
        [("Relational DB", GREEN), ("Object storage", GREEN), ("Vector index", GREEN)],
        [("Audit service", AMBER), ("Model registry", AMBER), ("Logs / metrics / traces", AMBER)],
    ]
    for row_index, row in enumerate(rows):
        y = 2.05 + row_index * 1.05
        for col_index, (label, color) in enumerate(row):
            x = 0.9 + col_index * 4.15
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.55), Inches(0.72))
            shape.fill.solid()
            shape.fill.fore_color.rgb = WHITE
            shape.line.color.rgb = color
            add_text(slide, label, x + 0.1, y + 0.22, 3.35, 0.25, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, 1.15, 6.35, 11.0, 0.55, "Governance plane", "Reference approvals • model validation • security policy • retention • change control", accent=GREEN, title_size=9, body_size=8)
    add_footer(slide, 13)

    # 14. Roadmap
    slide = new_slide(presentation)
    add_title(slide, "Delivery roadmap", "Move from workflow validation to controlled production in measurable stages", "Roadmap")
    phases = [
        ("Now", "Prototype", "Workflow complete\nImage RAG\nExplainability\nHuman review", BLUE),
        ("0–3 months", "Controlled pilot", "Persistent data\nAuth and roles\nLabeled evaluation set\nPilot baseline", GREEN),
        ("3–6 months", "Production foundation", "Validated model\nAudit trail\nReference governance\nMonitoring", AMBER),
        ("6–12 months", "Scale", "Multiple lines/SKUs\nMES/QMS integration\nMLOps and drift\nHA / DR", NAVY),
    ]
    x = 0.55
    for period, title, body, color in phases:
        add_card(slide, x, 2.25, 2.9, 3.75, title, body, accent=color, title_size=15, body_size=10.5)
        add_text(slide, period, x + 0.15, 6.18, 2.6, 0.28, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        x += 3.15
    add_footer(slide, 14)

    # 15. Decisions
    slide = new_slide(presentation, DEEP_NAVY)
    add_text(slide, "DECISIONS AND NEXT STEPS", 0.7, 0.48, 4.5, 0.3, size=10, color=GREEN, bold=True)
    add_text(slide, "What must be agreed\nbefore the pilot starts?", 0.7, 1.15, 6.0, 1.2, size=30, color=WHITE, bold=True)
    decisions = [
        "Intended use and decision authority",
        "Pilot product, SKU, line, and camera setup",
        "False-accept and false-reject tolerances",
        "Reference-image ownership and approval",
        "Labeled defect data and taxonomy",
        "Security, retention, audit, and validation scope",
    ]
    add_bullets(slide, decisions, 0.78, 2.75, 6.2, 3.35, size=14, color=RGBColor(224, 242, 254))
    add_card(slide, 7.5, 1.35, 4.95, 4.8, "Recommended immediate action", "Approve a controlled pilot charter.\n\nSelect one packaging family and imaging station.\n\nCollect adjudicated good and defective examples.\n\nMeasure accuracy, review time, and override rates.\n\nUse results to approve—or stop—the production investment.", fill=WHITE, accent=GREEN, title_size=16, body_size=12)
    add_text(slide, "PharmaInspect AI", 0.75, 6.95, 3.0, 0.25, size=9, color=MUTED, bold=True)

    presentation.save(PPTX_OUTPUT)


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print(DOCX_OUTPUT)
    print(PPTX_OUTPUT)
